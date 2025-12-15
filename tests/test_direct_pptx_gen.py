import sys
import os
sys.path.append(os.getcwd())
from src.publishing.pptx_generator import PptxGenerator
from pptx import Presentation

def test_generator():
    # 1. Create Dummy Template
    template_path = "dummy_template.pptx"
    p = Presentation()
    # Add a master/layout
    master = p.slide_masters[0]
    # Ensure it has layouts
    if not master.slide_layouts:
         p.save(template_path)
    else:
         p.save(template_path)

    # 2. Config
    config = {
        "mappings": {
            "documentary": {
                "master_index": 0,
                "layout_index": 0,
                "placeholders": {"title": 0, "body": 1}
            }
        }
    }
    
    # 3. Nodes
    nodes = [
        {"title": "Test Slide 1", "content_markdown": "## Header\n* Bullet 1", "target_layout": "documentary"}
    ]
    
    # 4. Generate
    output_path = "output_test.pptx"
    gen = PptxGenerator(config, template_path)
    gen.generate("Test Project", nodes, output_path)
    
    # 5. Verify
    if os.path.exists(output_path):
        print("Success: PPTX generated.")
        prs = Presentation(output_path)
        print(f"Slides: {len(prs.slides)}")
        print(f"Slide 1 Title: {prs.slides[1].shapes.title.text if len(prs.slides) > 1 else 'N/A'}")
    else:
        print("Failure: No output.")

    # Cleanup
    try:
        os.remove(template_path)
        os.remove(output_path)
    except: pass

if __name__ == "__main__":
    test_generator()
