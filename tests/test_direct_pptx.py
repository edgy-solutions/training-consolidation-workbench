import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_shapes(pptx_path):
    print(f"\n--- Inspecting {pptx_path} ---")
    if not os.path.exists(pptx_path):
        print("File not found.")
        return

    prs = Presentation(pptx_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Slide {slide_idx + 1}:")
        for shape in slide.shapes:
            print(f"  - Shape: {shape.name}, Type: {shape.shape_type}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"    -> PICTURE detected. Format: {shape.image.ext}")
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                print(f"    -> PLACEHOLDER. Type: {shape.placeholder_format.type}")
                if hasattr(shape, "image"):
                     print(f"       -> Has image! Format: {shape.image.ext}")
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print("    -> GROUP")
                for subshape in shape.shapes:
                    print(f"      - SubShape: {subshape.name}, Type: {subshape.shape_type}")

if __name__ == "__main__":
    files = [
        r"test_docs\course_b_git_workflows.pptx",
        r"test_docs\course_e_thermodynamics.pptx"
    ]
    for f in files:
        inspect_shapes(f)
