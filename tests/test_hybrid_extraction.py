import os
import shutil
import tempfile
from pptx import Presentation
from pptx.util import Inches
from src.ingestion.pptx_media_extractor import extract_images_from_pptx
from src.ingestion.rendering import convert_to_pdf
from src.ingestion.extraction import extract_text_and_metadata

def create_test_pptx_with_image(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
    
    # Create a dummy image
    from PIL import Image
    img = Image.new('RGB', (100, 100), color = 'red')
    img_path = "temp_test_image.png"
    img.save(img_path)
    
    # Add to slide
    slide.shapes.add_picture(img_path, Inches(1), Inches(1))
    
    prs.save(filename)
    os.remove(img_path)
    print(f"Created {filename} with a standard image.")

def test_hybrid_logic(pptx_path):
    print(f"\nTesting hybrid logic on: {pptx_path}")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_extract_dir = os.path.join(temp_dir, "extracted")
        os.makedirs(temp_extract_dir)
        
        # 1. Try Direct Extraction
        print("1. Attempting Direct Extraction...")
        direct_images = extract_images_from_pptx(pptx_path, temp_extract_dir)
        print(f"   Direct extraction found: {len(direct_images)} images")
        
        all_images = [f for f in os.listdir(temp_extract_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        
        # 2. Fallback
        if not all_images:
            print("   No images found. Triggering PDF Fallback...")
            try:
                pdf_path = convert_to_pdf(pptx_path, temp_dir)
                print(f"   Converted to PDF: {pdf_path}")
                
                elements = extract_text_and_metadata(
                    pdf_path,
                    extract_images=True,
                    image_output_dir=temp_extract_dir
                )
                
                pdf_images = [f for f in os.listdir(temp_extract_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
                print(f"   PDF Fallback found: {len(pdf_images)} images (total in dir)")
            except Exception as e:
                print(f"   PDF Fallback failed: {e}")
        else:
            print("   Direct extraction successful. Skipping fallback.")

if __name__ == "__main__":
    # Test Case A: Standard Image (Should succeed with Direct)
    pptx_a = "test_standard.pptx"
    create_test_pptx_with_image(pptx_a)
    test_hybrid_logic(pptx_a)
    
    # Test Case B: The problematic sample (Should trigger Fallback)
    pptx_b = r"test_docs\course_b_git_workflows.pptx"
    if os.path.exists(pptx_b):
        test_hybrid_logic(pptx_b)
    else:
        print(f"Skipping Test Case B: {pptx_b} not found")
        
    # Cleanup
    if os.path.exists(pptx_a):
        os.remove(pptx_a)
