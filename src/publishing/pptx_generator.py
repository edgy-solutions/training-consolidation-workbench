from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import json
import requests
import tempfile
import os
from typing import List, Dict, Any, Tuple
from PIL import Image as PILImage

# --- Helper Functions (Static) ---

def download_image(url: str, temp_dir: str) -> Tuple[str | None, Tuple[int, int] | None]:
    """
    Downloads an image from a URL and returns (local file path, (width, height)).
    """
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        
        content_type = response.headers.get('content-type', '')
        if 'png' in content_type or '.png' in url.lower():
            ext = '.png'
        elif 'gif' in content_type or '.gif' in url.lower():
            ext = '.gif'
        else:
            ext = '.jpg'
        
        filename = f"image_{hash(url) % 10000}{ext}"
        filepath = os.path.join(temp_dir, filename)
        
        with open(filepath, 'wb') as f:
            f.write(response.content)
        
        try:
            with PILImage.open(filepath) as img:
                dimensions = img.size
        except:
            dimensions = (400, 300)
        
        return filepath, dimensions
    except Exception as e:
        print(f"Failed to download image from {url}: {e}")
        return None, None

def parse_content_segments(markdown_text: str) -> List[Tuple[str, Any]]:
    """
    Parses markdown and returns a list of (type, content) tuples.
    """
    if not markdown_text:
        return [('text', '')]
    
    image_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
    result = []
    last_end = 0
    
    for match in re.finditer(image_pattern, markdown_text):
        text_before = markdown_text[last_end:match.start()].strip()
        if text_before:
            result.append(('text', text_before))
        
        alt_text = match.group(1)
        size_info = None
        filename = alt_text
        
        if '|' in alt_text and '{' in alt_text:
            parts = alt_text.split('|', 1)
            filename = parts[0]
            try:
                size_info = json.loads(parts[1])
            except json.JSONDecodeError:
                pass
        
        result.append(('image', {
            'alt': filename,
            'url': match.group(2),
            'size_info': size_info
        }))
        
        last_end = match.end()
    
    remaining_text = markdown_text[last_end:].strip()
    if remaining_text:
        result.append(('text', remaining_text))
    
    return result if result else [('text', '')]

def add_text_to_frame(text_frame, markdown_text):
    """
    Parses simple markdown (bullets, bold, headers) into a PPTX text frame.
    """
    if not markdown_text:
        return

    lines = markdown_text.split('\n')
    first_para = True
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        is_header = False
        if line.startswith('#'):
            line = re.sub(r'^#+\s*', '', line)
            is_header = True
            
        level = 0
        clean_line = line
        
        if line.startswith('- ') or line.startswith('* '):
            level = 0
            clean_line = line[2:]
        elif line.startswith('  - ') or line.startswith('  * '):
            level = 1
            clean_line = line[4:]
        
        if first_para:
            try:
                p = text_frame.paragraphs[0]
                first_para = False
            except IndexError:
                p = text_frame.add_paragraph()
        else:
            p = text_frame.add_paragraph()
        p.level = level
        
        parts = re.split(r'(\*\*.*?\*\*)', clean_line)
        
        for part in parts:
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run.text = part
            run.font.size = Pt(11)
            if is_header:
                run.font.bold = True
                run.font.size = Pt(14)
                
        p.space_after = Pt(6)

# --- Main Class ---

class PptxGenerator:
    def __init__(self, config: Dict[str, Any], template_file_path: str = None):
        """
        Initialize generator with configuration and a local path to the PPTX template.
        """
        self.config = config
        self.temp_dir = tempfile.mkdtemp()
        
        if template_file_path and os.path.exists(template_file_path):
            self.prs = Presentation(template_file_path)
            self._clean_template_slides()
        else:
            print("Warning: No valid template path provided. Using blank presentation.")
            self.prs = Presentation()

    def _clean_template_slides(self):
        """Removes all slides from the template to start fresh."""
        try:
            for i in range(len(self.prs.slides) - 1, -1, -1):
                rId = self.prs.slides._sldIdLst[i].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[i]
        except Exception as e:
            print(f"Warning: Failed to clear template slides: {e}")

    def _get_layout_config(self, archetype: str) -> Dict[str, Any]:
        """
        Retrieves configuration for a specific archetype from the mappings.
        Falls back to 'documentary' or the first available if not found.
        """
        mappings = self.config.get('mappings', {})
        if archetype in mappings:
            return mappings[archetype]
        
        # Fallback to documentary
        if 'documentary' in mappings:
            return mappings['documentary']
            
        # Ultimate fallback (first one)
        if mappings:
            return list(mappings.values())[0]
            
        # No mappings defined? Return decent default
        return {
            "master_index": 0,
            "layout_index": 1,
            "placeholders": {"title": 0, "body": 1}
        }

    def _get_placeholder(self, slide, name: str, ph_map: Dict[str, int]):
        """Safely retrieve a placeholder by its role name using the mapping."""
        idx = ph_map.get(name)
        if idx is None:
            return None
        try:
            return slide.placeholders[idx]
        except KeyError:
            return None

    def add_slide(self, archetype: str) -> Any:
        """
        Creates a new slide based on the archetype.
        Attaches the placeholder mapping to the slide object for later use.
        """
        layout_cfg = self._get_layout_config(archetype)
        m_idx = layout_cfg.get('master_index', 0)
        l_idx = layout_cfg.get('layout_index', 0)
        
        # Get Layout
        try:
            if m_idx < len(self.prs.slide_masters):
                master = self.prs.slide_masters[m_idx]
                if l_idx < len(master.slide_layouts):
                    layout = master.slide_layouts[l_idx]
                else:
                    layout = master.slide_layouts[0]
            else:
                layout = self.prs.slide_layouts[0]
        except:
             layout = self.prs.slide_layouts[0]
             
        slide = self.prs.slides.add_slide(layout)
        
        # Attach configuration to slide object (runtime monkey-patching for convenience)
        slide.ph_map = layout_cfg.get('placeholders', {})
        return slide

    def insert_content(self, slide, content_segments):
        """
        Intelligently distributes content segments into the available placeholders 
        defined in the slide's configuration.
        """
        ph_map = getattr(slide, 'ph_map', {})
        
        # Get all text and image segments
        text_segments = [s for s in content_segments if s[0] == 'text']
        image_segments = [s for s in content_segments if s[0] == 'image']
        
        # Strategy depends on available placeholders
        # We try to fill specific roles first
        
        # 1. Images
        # Try specific image roles: 'image', 'col1_img', 'col2_img', etc.
        # Or generic 'right_body' if split
        
        image_roles = ['image', 'col1_img', 'col2_img', 'col3_img', 'right_body']
        img_idx = 0
        
        for role in image_roles:
            if role in ph_map and img_idx < len(image_segments):
                ph = self._get_placeholder(slide, role, ph_map)
                if ph:
                    try:
                        url = image_segments[img_idx][1]['url']
                        img_path, _ = download_image(url, self.temp_dir)
                        if img_path:
                            ph.insert_picture(img_path)
                            img_idx += 1
                    except Exception as e:
                        print(f"Failed to insert image into {role}: {e}")

        # 2. Text
        # Try specific text roles: 'body', 'left_body', 'col1_body', etc.
        text_roles = ['body', 'left_body', 'col1_body', 'col2_body', 'col3_body']
        txt_idx = 0
        
        for role in text_roles:
            if role in ph_map and txt_idx < len(text_segments):
                ph = self._get_placeholder(slide, role, ph_map)
                if ph:
                    if not ph.has_text_frame: continue
                    add_text_to_frame(ph.text_frame, text_segments[txt_idx][1])
                    txt_idx += 1
                    
        # Fallback: If we have leftover content, append to main body if possible
        if txt_idx < len(text_segments):
            # Find the first valid body placeholder
            main_body_role = next((r for r in ['body', 'left_body'] if r in ph_map), None)
            if main_body_role:
                ph = self._get_placeholder(slide, main_body_role, ph_map)
                if ph and ph.has_text_frame:
                    remaining_text = "\n\n".join([s[1] for s in text_segments[txt_idx:]])
                    # Add separator
                    p = ph.text_frame.add_paragraph()
                    p.text = "---"
                    add_text_to_frame(ph.text_frame, remaining_text)

    def generate(self, project_title: str, nodes: List[Dict[str, Any]], output_path: str):
        """
        Main execution method.
        """
        # 1. Title Slide
        slide = self.add_slide('hero')
        
        title_ph = self._get_placeholder(slide, 'title', slide.ph_map)
        if title_ph: 
            title_ph.text = project_title
        elif slide.shapes.title:
            slide.shapes.title.text = project_title
            
        sub_ph = self._get_placeholder(slide, 'subtitle', slide.ph_map)
        if sub_ph:
            sub_ph.text = f"Generated by Training Consolidation Workbench\n{self.config.get('template_name', '')}"
            
        # 2. Content Slides
        sorted_nodes = sorted(nodes, key=lambda x: x.get('order', 0))
        
        for node in sorted_nodes:
            target_layout = node.get('target_layout', 'documentary')
            slide = self.add_slide(target_layout)
            
            # Set Title
            title_ph = self._get_placeholder(slide, 'title', slide.ph_map)
            if title_ph:
                title_ph.text = node.get('title', 'Untitled')
            elif slide.shapes.title:
                slide.shapes.title.text = node.get('title', 'Untitled')
                
            # Content
            content_md = node.get('content_markdown', '')
            segments = parse_content_segments(content_md)
            
            self.insert_content(slide, segments)
            
        # Save
        self.prs.save(output_path)
        
        # Cleanup
        try:
            import shutil
            shutil.rmtree(self.temp_dir)
        except: pass
        
        return output_path
