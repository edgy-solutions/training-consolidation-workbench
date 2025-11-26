import os
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches

def generate_course_a_pdf(output_dir):
    filename = "course_a_intro_to_git.pdf"
    path = os.path.join(output_dir, filename)
    
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    
    # Page 1: Title
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 100, "Course A: Intro to Git")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, "Module 1: Git Basics")
    
    # Content
    c.drawString(100, height - 160, "Git is a distributed version control system.")
    c.drawString(100, height - 180, "Key concepts: Repository, Commit, Branch.")
    
    # Diagram: Commit Tree
    c.setStrokeColorRGB(0, 0, 0)
    c.setLineWidth(2)
    
    # Node 1
    c.circle(150, height - 300, 20, fill=0)
    c.drawString(140, height - 300, "C1")
    
    # Node 2
    c.circle(250, height - 300, 20, fill=0)
    c.drawString(240, height - 300, "C2")
    
    # Edge
    c.line(170, height - 300, 230, height - 300)
    
    c.drawString(100, height - 350, "Figure 1: A simple linear commit history.")
    
    c.showPage()
    c.save()
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Ray",
        "course_title": "Intro to Git",
        "version": "1.0",
        "scope_of_material": "Git Basics, Commits, Branches",
        "current_delivery_method": "Web",
        "duration_hours": 2.0,
        "costs": "Low",
        "tech_data_assessment": "None",
        "source_of_content": "Internal Wiki",
        "current_instructors": "Alice Smith",
        "audience": "Junior Developers",
        "location": "Remote",
        "level_of_material": "Beginner",
        "engineering_discipline": "Software",
        "comments": "Standard onboarding course."
    }
    return filename, path, metadata

def generate_course_b_pptx(output_dir):
    filename = "course_b_git_workflows.pptx"
    path = os.path.join(output_dir, filename)
    
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Course B: Git Workflows"
    subtitle.text = "Topic: Feature Branches and PRs"
    
    # Slide 2: Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Feature Branch Workflow"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "1. Create a branch from main"
    p = tf.add_paragraph()
    p.text = "2. Commit changes"
    p = tf.add_paragraph()
    p.text = "3. Open a Pull Request"
    
    # Slide 3: Diagram (Shapes)
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Merge Flowchart"
    
    # Main branch line
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(6), Inches(0.2))
    slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(0.5)).text_frame.text = "Main Branch"
    
    # Feature branch arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.CURVED_UP_ARROW, Inches(2), Inches(3.5), Inches(3), Inches(1))
    arrow.text_frame.text = "Feature Branch"
    
    prs.save(path)
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Collin",
        "course_title": "Git Workflows",
        "version": "2.1",
        "scope_of_material": "Feature Branching, Pull Requests",
        "current_delivery_method": "In Person",
        "duration_hours": 4.0,
        "costs": "Medium",
        "tech_data_assessment": "Required",
        "source_of_content": "DevOps Handbook",
        "current_instructors": "Bob Jones",
        "audience": "All Developers",
        "location": "New York Office",
        "level_of_material": "Intermediate",
        "engineering_discipline": "Software",
        "comments": "Critical for compliance."
    }
    return filename, path, metadata

def generate_course_c_docx(output_dir):
    filename = "course_c_git_collaboration.docx"
    path = os.path.join(output_dir, filename)
    
    doc = Document()
    doc.add_heading('Course C: Collaborating with Git', 0)
    
    doc.add_heading('Resolving Conflicts', level=1)
    doc.add_paragraph('Conflicts occur when two branches modify the same line of code.')
    
    doc.add_heading('Merge vs Rebase', level=1)
    table = doc.add_table(rows=3, cols=2)
    table.style = 'Table Grid'
    
    # Header
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Strategy'
    hdr_cells[1].text = 'Pros/Cons'
    
    # Row 1
    row1 = table.rows[1].cells
    row1[0].text = 'Merge'
    row1[1].text = 'Preserves history, but can be messy.'
    
    # Row 2
    row2 = table.rows[2].cells
    row2[0].text = 'Rebase'
    row2[1].text = 'Clean history, but rewrites commits.'
    
    doc.save(path)
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Pat",
        "course_title": "Collaborating with Git",
        "version": "1.5",
        "scope_of_material": "Conflicts, Merge vs Rebase",
        "current_delivery_method": "Hybrid",
        "duration_hours": 3.0,
        "costs": "Low",
        "tech_data_assessment": "None",
        "source_of_content": "Community Guidelines",
        "current_instructors": "Charlie Brown",
        "audience": "Senior Developers",
        "location": "London Office",
        "level_of_material": "Advanced",
        "engineering_discipline": "Software",
        "comments": "Focus on conflict resolution."
    }
    return filename, path, metadata

def generate_all(output_dir="test_docs"):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    docs = []
    docs.append(generate_course_a_pdf(output_dir))
    docs.append(generate_course_b_pptx(output_dir))
    docs.append(generate_course_c_docx(output_dir))
    return docs

if __name__ == "__main__":
    generate_all()
