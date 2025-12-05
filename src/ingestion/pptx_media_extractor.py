import os
import zipfile
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def _iter_picture_shapes(prs):
    """
    Iterate through all picture shapes in the presentation, preserving slide index.
    Yields (slide_index, shape).
    """
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield slide_idx, shape
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    if subshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        yield slide_idx, subshape
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "image"):
                    yield slide_idx, shape

def extract_images_from_pptx(pptx_path: str, output_dir: str) -> List[str]:
    """
    Extract images from a PPTX file using python-pptx.
    Returns a list of paths to the extracted images.
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
        
    extracted_files = []
    
    try:
        prs = Presentation(pptx_path)
        
        # Track seen images to avoid duplicates if needed, or just save all
        # We'll save all to preserve context
        
        count = 0
        for slide_idx, shape in _iter_picture_shapes(prs):
            try:
                image = shape.image
                ext = image.ext
                # Use a naming convention that preserves slide order
                # slide_{i}_img_{j}.ext
                filename = f"slide_{slide_idx+1}_img_{count}.{ext}"
                filepath = os.path.join(output_dir, filename)
                
                with open(filepath, 'wb') as f:
                    f.write(image.blob)
                
                extracted_files.append(filepath)
                count += 1
            except Exception as e:
                print(f"Warning: Failed to extract image from shape on slide {slide_idx+1}: {e}")
                continue
                
    except Exception as e:
        print(f"Error reading PPTX {pptx_path}: {e}")
        # We don't raise here, we just return what we got (or empty list)
        # so the caller can decide to use fallback
        pass
        
    return extracted_files

def extract_media_via_zip(pptx_path: str, output_dir: str) -> List[str]:
    """
    Extract all media (images, video, audio) by treating PPTX as a ZIP.
    Useful for getting original SVGs or videos that python-pptx might miss.
    Note: Loses slide context.
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
        
    extracted_files = []
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            for file_info in z.infolist():
                if file_info.filename.startswith('ppt/media/'):
                    # Extract to output_dir, flattening the path
                    filename = os.path.basename(file_info.filename)
                    target_path = os.path.join(output_dir, f"zip_extracted_{filename}")
                    
                    with open(target_path, 'wb') as f:
                        f.write(z.read(file_info))
                    
                    extracted_files.append(target_path)
    except Exception as e:
        print(f"Error zip-extracting PPTX {pptx_path}: {e}")
        
    return extracted_files
